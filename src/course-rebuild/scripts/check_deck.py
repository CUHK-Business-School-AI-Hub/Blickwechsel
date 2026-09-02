#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
check_deck.py - look for the layout damage you cannot see without opening a deck.

Run this after any script edits a .pptx. Nobody can view a rendered slide from
a script, but a lot of what goes wrong is geometry, and geometry is in the file.

  REFUSE   a shape runs off the edge of its slide
  WARN     an empty placeholder, a table too wide or too tall, tiny type,
           a slide with no title, or far more text than its box will hold
  NOTE     counts, so you can see at a glance what changed

Exits non-zero if anything is in the REFUSE class, so it can sit in front of a
publish step.

Usage:
    python check_deck.py "deck.pptx"
    python check_deck.py "new.pptx" --against "original.pptx"   # compare
    python check_deck.py "deck.pptx" --slides 5,6,11            # only these

Needs python-pptx:  py -m pip install python-pptx
"""

import argparse
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("Needs python-pptx.  Install it with:  py -m pip install python-pptx")

IN = 914400.0
MIN_PT = 10          # smaller than this is unreadable from the back of a room
MAX_COLS = 6         # a wider table is a spreadsheet, not a slide
CHARS_PER_SQIN = 44  # rough: above this the box is over-full


def audit(path, only=None):
    p = Presentation(path)
    W, H = p.slide_width, p.slide_height
    refuse, warn, note = [], [], []
    n_slides = 0

    for i, s in enumerate(p.slides, start=1):
        n_slides = i
        if only and i not in only:
            continue
        has_title, texts = False, 0

        for sh in s.shapes:
            name = sh.name[:24]
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                l = t = w = h = None

            if None not in (l, t, w, h):
                over = []
                if l + w > W:
                    over.append("%.2f in past the right edge" % ((l + w - W) / IN))
                if t + h > H:
                    over.append("%.2f in past the bottom" % ((t + h - H) / IN))
                if l < 0:
                    over.append("%.2f in off the left" % (-l / IN))
                if t < 0:
                    over.append("%.2f in off the top" % (-t / IN))
                if over:
                    refuse.append((i, name, "; ".join(over)))

            if sh.has_text_frame:
                txt = sh.text_frame.text
                if txt.strip():
                    texts += 1
                if sh.is_placeholder:
                    try:
                        idx = sh.placeholder_format.idx
                    except Exception:
                        idx = None
                    if idx == 0 and txt.strip():
                        has_title = True
                    if not txt.strip():
                        warn.append((i, name, "placeholder left empty"))
                if txt.strip() and None not in (w, h) and w and h:
                    area = (w / IN) * (h / IN)
                    if area > 0 and len(txt) / area > CHARS_PER_SQIN:
                        warn.append((i, name, "%d characters in %.1f square inches "
                                              "- likely to overflow its box"
                                     % (len(txt), area)))
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        sz = run.font.size
                        if sz is not None and sz.pt < MIN_PT:
                            warn.append((i, name, "type at %.0f pt" % sz.pt))
                            break
                    else:
                        continue
                    break

            if sh.has_table:
                tbl = sh.table
                cols = len(tbl.columns)
                if cols > MAX_COLS:
                    warn.append((i, name, "table with %d columns" % cols))
                rows = len(tbl.rows)
                if None not in (t, h) and t + h > H * 0.95:
                    warn.append((i, name, "table of %d rows reaches the bottom "
                                          "of the slide" % rows))

        if not has_title:
            note.append((i, "-", "no title placeholder with text on it"))

    return {"path": path, "slides": n_slides, "w": W / IN, "h": H / IN,
            "refuse": refuse, "warn": warn, "note": note}


def show(label, rows):
    if not rows:
        return
    print("  %s" % label)
    for i, name, msg in rows:
        print("     slide %-3d %-24s %s" % (i, name, msg))
    print()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--against", default=None,
                    help="the original, to compare slide counts against")
    ap.add_argument("--slides", default=None,
                    help="only check these, e.g. 5,6,11")
    args = ap.parse_args()

    only = None
    if args.slides:
        only = set(int(x) for x in args.slides.replace(" ", "").split(",") if x)

    r = audit(args.pptx, only)
    W = 70
    print("=" * W)
    print("DECK CHECK  -  %s" % r["path"].split("\\")[-1].split("/")[-1])
    print("=" * W)
    print()
    print("  %d slides, %.2f x %.2f inches" % (r["slides"], r["w"], r["h"]))
    if only:
        print("  checking only slides %s" % ", ".join(str(x) for x in sorted(only)))
    if args.against:
        o = audit(args.against)
        print("  original had %d slides; this has %d  (%+d)"
              % (o["slides"], r["slides"], r["slides"] - o["slides"]))
    print()

    refuse, warn, inherited = r["refuse"], r["warn"], []
    if args.against:
        # A problem the original already had is not something this edit broke.
        # Report it, but do not refuse on it - a gate that blocks on somebody
        # else's pre-existing mess is a gate people learn to ignore.
        o = audit(args.against)
        was = set((name, msg) for _, name, msg in o["refuse"] + o["warn"])
        inherited = [x for x in refuse + warn if (x[1], x[2]) in was]
        refuse = [x for x in refuse if (x[1], x[2]) not in was]
        warn = [x for x in warn if (x[1], x[2]) not in was]

    show("RUNS OFF THE SLIDE - this edit broke it, fix before anyone sees it:", refuse)
    show("WORTH OPENING AND LOOKING AT:", warn)
    if inherited:
        show("ALREADY LIKE THIS IN THE ORIGINAL - not caused by this edit:", inherited)
    show("NOTES:", r["note"])

    if not refuse and not warn:
        print("  Nothing this edit changed runs off a slide or looks over-full.")
        print()

    print("-" * W)
    print("This checks geometry, not whether a slide READS well. Line breaks in")
    print("awkward places, a table that fits but nobody can follow, a colour that")
    print("vanishes on a projector - none of that is in the file. Open the deck.")
    return 1 if refuse else 0


if __name__ == "__main__":
    sys.exit(main())
